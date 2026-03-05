import os
import re
import uuid
import base64
import logging
import tempfile
import traceback
import json
import requests
from io import BytesIO
from datetime import datetime, timedelta, timezone
from agents.basic_agent import BasicAgent
from utils.azure_file_storage import AzureFileStorageManager
# Always import python-pptx modules unconditionally
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from azure.storage.file import FileService


class PowerPointAgent(BasicAgent):
    """
    PowerPoint Agent using Microsoft Graph API for creating presentations.
    Features enhanced customization capabilities for slide content and design.

    Note: Requires Microsoft Graph API credentials to be set in environment variables:
    - GRAPH_CLIENT_ID
    - GRAPH_CLIENT_SECRET
    - GRAPH_TENANT_ID
    """

    def __init__(self):
        self.name = 'PowerPoint'
        self.metadata = {
            "name": self.name,
            "description": "Creates professional PowerPoint presentations with extensive customization options.",
            "parameters": {
                "type": "object",
                "properties": {
                    # Presentation basics
                    "title": {
                        "type": "string",
                        "description": "The main title of your presentation. This will appear on the title slide and in the filename."
                    },
                    "subtitle": {
                        "type": "string",
                        "description": "Optional subtitle for the title slide. If not provided, today's date will be used."
                    },
                    "author": {
                        "type": "string",
                        "description": "Optional author name to include on the title slide."
                    },

                    # Theming and styling
                    "theme": {
                        "type": "string",
                        "description": "Presentation theme to use. Options: 'professional' (blue/gray corporate style), 'creative' (vibrant colors for engaging presentations), 'minimalist' (clean, simple design), 'academic' (formal style for research presentations), 'custom' (use provided colors).",
                        "enum": ["professional", "creative", "minimalist", "academic", "custom"],
                        "default": "professional"
                    },
                    "primary_color": {
                        "type": "string",
                        "description": "Primary color for the presentation theme in hex format (e.g., '#0072C6'). Used when theme is set to 'custom'."
                    },
                    "secondary_color": {
                        "type": "string",
                        "description": "Secondary color for the presentation theme in hex format (e.g., '#00B294'). Used when theme is set to 'custom'."
                    },
                    "font": {
                        "type": "string",
                        "description": "Main font to use throughout the presentation. Options: 'calibri', 'arial', 'times', 'georgia', 'verdana', 'tahoma'.",
                        "enum": ["calibri", "arial", "times", "georgia", "verdana", "tahoma"],
                        "default": "calibri"
                    },

                    # Title slide
                    "include_title_slide": {
                        "type": "boolean",
                        "description": "Whether to include a title slide at the beginning of the presentation.",
                        "default": True
                    },
                    "title_slide_image": {
                        "type": "string",
                        "description": "Optional placeholder for an image on the title slide. Format as 'image:width:height' (e.g., 'logo:2:1' for a logo 2 inches wide and 1 inch tall)."
                    },

                    # Agenda slide
                    "include_agenda": {
                        "type": "boolean",
                        "description": "Whether to include an agenda/table of contents slide after the title slide.",
                        "default": False
                    },
                    "agenda_title": {
                        "type": "string",
                        "description": "Title for the agenda slide if included. Defaults to 'Agenda' if not specified.",
                        "default": "Agenda"
                    },
                    "auto_generate_agenda": {
                        "type": "boolean",
                        "description": "Whether to automatically generate agenda items from slide titles. If false, use 'agenda_items' to specify custom agenda.",
                        "default": True
                    },
                    "agenda_items": {
                        "type": "array",
                        "description": "Custom agenda items to list on the agenda slide if auto_generate_agenda is false.",
                        "items": {
                            "type": "string"
                        }
                    },

                    # Content slides
                    "slides": {
                        "type": "array",
                        "description": "Array of slide objects defining the main content slides.",
                        "items": {
                            "type": "object",
                            "properties": {
                                "title": {
                                    "type": "string",
                                    "description": "The title of the slide."
                                },
                                "subtitle": {
                                    "type": "string",
                                    "description": "Optional subtitle for the slide."
                                },
                                "content": {
                                    "type": "array",
                                    "description": "Array of content items for the slide. For bullet points, prefix with '-' or '*'. For numbered lists, prefix with '1.', '2.', etc.",
                                    "items": {
                                        "type": "string"
                                    }
                                },
                                "notes": {
                                    "type": "string",
                                    "description": "Optional speaker notes to add to the slide."
                                },
                                "layout": {
                                    "type": "string",
                                    "description": "The layout for the slide: 'title_and_content' (standard layout), 'section' (section divider), 'two_content' (side-by-side content), 'title_only' (just a title), 'blank' (empty slide), 'comparison' (comparison with left/right sides), 'quote' (highlighted quote layout).",
                                    "enum": [
                                        "title_and_content",
                                        "section",
                                        "two_content",
                                        "title_only",
                                        "blank",
                                        "comparison",
                                        "quote"
                                    ],
                                    "default": "title_and_content"
                                },
                                "content_left": {
                                    "type": "array",
                                    "description": "Content for the left side when using 'two_content' or 'comparison' layout.",
                                    "items": {
                                        "type": "string"
                                    }
                                },
                                "content_right": {
                                    "type": "array",
                                    "description": "Content for the right side when using 'two_content' or 'comparison' layout.",
                                    "items": {
                                        "type": "string"
                                    }
                                },
                                "image_placeholder": {
                                    "type": "string",
                                    "description": "Optional placeholder for an image on the slide. Format as 'image:width:height:position' (e.g., 'chart:4:3:center' for a chart 4 inches wide, 3 inches tall, centered). Position can be 'center', 'left', 'right', 'top', 'bottom'."
                                },
                                "background_color": {
                                    "type": "string",
                                    "description": "Optional hex color code for slide background, overriding the theme (e.g., '#F5F5F5'). Leave empty to use theme default."
                                }
                            },
                            "required": ["title"]
                        }
                    },

                    # Special slides
                    "section_slides": {
                        "type": "array",
                        "description": "Optional array of section divider slides to insert at specific positions. Each object should include 'title' and 'position' (index to insert at).",
                        "items": {
                            "type": "object",
                            "properties": {
                                "title": {
                                    "type": "string",
                                    "description": "The title of the section slide."
                                },
                                "position": {
                                    "type": "integer",
                                    "description": "The position (0-based index) where the section slide should be inserted."
                                },
                                "background_color": {
                                    "type": "string",
                                    "description": "Optional hex color code for section slide background (e.g., '#0072C6'). Leave empty to use theme default."
                                }
                            },
                            "required": ["title", "position"]
                        }
                    },
                    "thank_you_slide": {
                        "type": "boolean",
                        "description": "Whether to include a 'Thank You' slide at the end of the presentation.",
                        "default": False
                    },
                    "thank_you_title": {
                        "type": "string",
                        "description": "Title for the thank you slide. Defaults to 'Thank You' if not specified.",
                        "default": "Thank You"
                    },
                    "thank_you_content": {
                        "type": "string",
                        "description": "Optional content for the thank you slide (e.g., contact information, next steps)."
                    },

                    # Output options
                    "output_format": {
                        "type": "string",
                        "description": "Output format: 'azure_storage' (save to Azure File Storage), 'base64' (return as base64 string), or 'local_file' (save to local filesystem).",
                        "enum": ["azure_storage", "base64", "local_file"],
                        "default": "azure_storage"
                    },
                    "filename_prefix": {
                        "type": "string",
                        "description": "Optional prefix for the generated filename. The final filename will be prefix_title_timestamp.pptx."
                    },
                    "azure_directory": {
                        "type": "string",
                        "description": "Directory in Azure File Storage for saving the presentation when using 'azure_storage' output format.",
                        "default": "presentations"
                    },
                    "file_path": {
                        "type": "string",
                        "description": "File path for local saving when using 'local_file' output format. If not specified, a temporary directory will be used."
                    },
                    "user_guid": {
                        "type": "string",
                        "description": "Optional user GUID for user-specific storage when using 'azure_storage' output format."
                    },

                    # Download link options
                    "generate_download_link": {
                        "type": "boolean",
                        "description": "Whether to generate a temporary download link after creating the presentation.",
                        "default": False
                    },
                    "download_link_expiry": {
                        "type": "integer",
                        "description": "Number of minutes the download link should remain valid. Only used if generate_download_link is true.",
                        "default": 30
                    },
                    "technical_path": {
                        "type": "string",
                        "description": "The technical path to an existing file for generating download links. Used when only generating a link without creating a new presentation."
                    }
                },
                "required": ["title", "slides"]
            }
        }
        self.storage_manager = AzureFileStorageManager()
        self.access_token = None

        # Try to get authentication details
        self.client_id = os.environ.get('GRAPH_CLIENT_ID')
        self.client_secret = os.environ.get('GRAPH_CLIENT_SECRET')
        self.tenant_id = os.environ.get('GRAPH_TENANT_ID')

        # Check if we have Graph API credentials
        self.graph_api_available = all(
            [self.client_id, self.client_secret, self.tenant_id])
        if not self.graph_api_available:
            logging.warning(
                "Microsoft Graph API credentials not found. Will use python-pptx instead.")

        # Define theme colors
        self.themes = {
            "professional": {
                "primary": "0072C6",   # Microsoft blue
                "secondary": "2F5597",  # Darker blue
                "background": "FFFFFF",  # White
                "accent1": "5B9BD5",    # Light blue
                "accent2": "ED7D31",    # Orange
                "text": "333333"        # Dark gray
            },
            "creative": {
                "primary": "FF5722",    # Deep orange
                "secondary": "8BC34A",  # Light green
                "background": "FAFAFA",  # Off-white
                "accent1": "03A9F4",    # Light blue
                "accent2": "9C27B0",    # Purple
                "text": "424242"        # Dark gray
            },
            "minimalist": {
                "primary": "212121",    # Near black
                "secondary": "757575",  # Medium gray
                "background": "FFFFFF",  # White
                "accent1": "BDBDBD",    # Light gray
                "accent2": "9E9E9E",    # Gray
                "text": "212121"        # Near black
            },
            "academic": {
                "primary": "003366",    # Deep blue
                "secondary": "990000",  # Deep red
                "background": "FFFFFF",  # White
                "accent1": "336699",    # Medium blue
                "accent2": "993333",    # Medium red
                "text": "333333"        # Dark gray
            }
        }

        # Set Presentation attribute directly
        self.Presentation = Presentation

        super().__init__(name=self.name, metadata=self.metadata)

    def perform(self, **kwargs):
        """
        Main entry point for creating PowerPoint presentations.

        Args:
            Many parameters as defined in metadata

        Returns:
            dict: Results of the operation
        """
        try:
            # Check if we're only generating a download link for an existing file
            if kwargs.get('technical_path') and not kwargs.get('title') and not kwargs.get('slides'):
                # In this case, we're just generating a download link for an existing file
                return self.generate_download_link(**kwargs)

            # Extract basic parameters
            title = kwargs.get('title', 'Untitled Presentation')
            slides = kwargs.get('slides', [])
            output_format = kwargs.get('output_format', 'azure_storage')
            azure_directory = kwargs.get('azure_directory', 'presentations')
            file_path = kwargs.get('file_path', '')
            user_guid = kwargs.get('user_guid', '')
            generate_download_link = kwargs.get(
                'generate_download_link', False)
            download_link_expiry = kwargs.get('download_link_expiry', 30)

            # Validate essential parameters
            if not title:
                return {"status": "error", "message": "Title cannot be empty"}

            if not slides:
                return {"status": "error", "message": "At least one slide is required"}

            # Create the presentation
            if self.graph_api_available:
                # Try Graph API first
                try:
                    presentation_data = self._create_presentation_graph(kwargs)
                except Exception as e:
                    logging.warning(
                        f"Graph API failed, falling back to python-pptx: {str(e)}")
                    presentation_data = self._create_presentation_pptx(kwargs)
            else:
                # Use python-pptx directly
                presentation_data = self._create_presentation_pptx(kwargs)

            # Handle output format
            result = None
            if output_format.lower() == 'azure_storage':
                result = self._save_to_azure(
                    presentation_data, title, azure_directory, user_guid, kwargs.get('filename_prefix'))
            elif output_format.lower() == 'base64':
                result = self._get_as_base64(presentation_data, title)
            else:  # local_file
                result = self._save_to_local_file(
                    presentation_data, title, file_path, kwargs.get('filename_prefix'))

            # Generate download link if requested and we have a successful azure storage save
            if generate_download_link and result and result.get('status') == 'success' and output_format.lower() == 'azure_storage':
                technical_path = result.get('technical_path')
                if technical_path:
                    download_result = self.generate_download_link(
                        file_path=technical_path,
                        expiry_minutes=download_link_expiry
                    )

                    # Merge the download link information into the result
                    if download_result and download_result.get('status') == 'success':
                        result['download_url'] = download_result.get(
                            'download_url')
                        result['download_expiry'] = download_result.get(
                            'expiry_time')
                        result[
                            'message'] += f" A download link has been generated that will expire in {download_link_expiry} minutes."

            return result

        except Exception as e:
            logging.error(f"Error in PowerPointAgent: {str(e)}")
            logging.error(traceback.format_exc())
            return {"status": "error", "message": f"Failed to create PowerPoint: {str(e)}"}

    def generate_download_link(self, file_path=None, expiry_minutes=30, **kwargs):
        """
        Generates a temporary download link for a previously created PowerPoint file.

        Args:
            file_path (str): Path to the file in Azure storage. If None, will use the one from kwargs.
            expiry_minutes (int): Number of minutes the link should remain valid
            **kwargs: Additional parameters, including previously saved technical_path

        Returns:
            dict: Result with download URL if successful
        """
        try:
            # If file_path is not provided directly, try to get it from kwargs
            if not file_path:
                file_path = kwargs.get('technical_path')

            if not file_path:
                return {"status": "error", "message": "No file path provided for download link generation"}

            # Parse the file path to get directory and filename
            path_parts = file_path.split('/')
            filename = path_parts[-1]
            directory = '/'.join(path_parts[:-1])

            # Generate SAS token with expiration time
            expiry_time = datetime.utcnow() + timedelta(minutes=expiry_minutes)

            # Get the download URL with SAS token
            download_url = self.storage_manager.generate_download_url(
                directory,
                filename,
                expiry_time
            )

            if download_url:
                return {
                    "status": "success",
                    "message": f"Download link generated successfully. Valid for {expiry_minutes} minutes.",
                    "download_url": download_url,
                    "expiry_time": expiry_time.isoformat(),
                    "filename": filename
                }
            else:
                return {
                    "status": "error",
                    "message": "Failed to generate download link"
                }

        except Exception as e:
            logging.error(f"Error generating download link: {str(e)}")
            logging.error(traceback.format_exc())
            return {
                "status": "error",
                "message": f"Error generating download link: {str(e)}"
            }

    def _get_theme_colors(self, kwargs):
        """
        Gets the theme colors based on parameters.

        Args:
            kwargs: Full parameter dictionary

        Returns:
            dict: Theme color dictionary
        """
        theme_name = kwargs.get('theme', 'professional')

        # If custom theme is selected and colors are provided, create custom theme
        if theme_name == 'custom':
            primary = kwargs.get('primary_color', '#0072C6').lstrip('#')
            secondary = kwargs.get('secondary_color', '#2F5597').lstrip('#')

            return {
                "primary": primary,
                "secondary": secondary,
                "background": "FFFFFF",
                "accent1": primary,
                "accent2": secondary,
                "text": "333333"
            }

        # Otherwise use predefined theme
        return self.themes.get(theme_name, self.themes['professional'])

    def _hex_to_rgb(self, hex_color):
        """
        Converts hex color to RGB tuple.

        Args:
            hex_color: Hex color code (with or without #)

        Returns:
            tuple: (R, G, B) values
        """
        hex_color = hex_color.lstrip('#')
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

    def _create_presentation_graph(self, kwargs):
        """
        Creates a PowerPoint presentation using Microsoft Graph API.

        Args:
            kwargs: All presentation parameters

        Returns:
            bytes: The binary presentation data
        """
        try:
            # Get access token
            if not self.access_token:
                self.access_token = self._get_access_token()
                if not self.access_token:
                    raise Exception("Failed to get Graph API access token")

            # Create a Drive Item with the presentation
            headers = {
                'Authorization': f'Bearer {self.access_token}',
                'Content-Type': 'application/json'
            }

            # Extract basic presentation info
            title = kwargs.get('title', 'Untitled Presentation')

            # Create empty presentation in the root of OneDrive
            sanitized_title = self._sanitize_filename(title)
            file_name = f"{sanitized_title}.pptx"

            # Define the presentation content payload
            presentation_payload = {
                "name": file_name,
                "file": {},
                "@microsoft.graph.conflictBehavior": "rename"
            }

            # Create empty file
            create_response = requests.post(
                "https://graph.microsoft.com/v1.0/me/drive/root/children",
                headers=headers,
                json=presentation_payload
            )
            create_response.raise_for_status()
            file_info = create_response.json()
            file_id = file_info.get('id')

            # Now add slides
            self._add_slides_graph(file_id, kwargs)

            # Download the file
            download_headers = {
                'Authorization': f'Bearer {self.access_token}'
            }

            download_response = requests.get(
                f"https://graph.microsoft.com/v1.0/me/drive/items/{file_id}/content",
                headers=download_headers
            )
            download_response.raise_for_status()

            # Delete the file from OneDrive (cleanup)
            requests.delete(
                f"https://graph.microsoft.com/v1.0/me/drive/items/{file_id}",
                headers=headers
            )

            # Return binary content
            return download_response.content

        except Exception as e:
            logging.error(
                f"Error creating PowerPoint with Graph API: {str(e)}")
            logging.error(traceback.format_exc())
            # If Graph API fails, fall back to python-pptx
            logging.info("Falling back to python-pptx")
            raise

    def _add_slides_graph(self, file_id, kwargs):
        """
        Adds slides to a presentation using Graph API.

        Args:
            file_id (str): The file ID in OneDrive
            kwargs: All presentation parameters
        """
        try:
            headers = {
                'Authorization': f'Bearer {self.access_token}',
                'Content-Type': 'application/json'
            }

            # Extract necessary parameters
            title = kwargs.get('title', 'Untitled Presentation')
            subtitle = kwargs.get(
                'subtitle', f"Created on {datetime.now().strftime('%B %d, %Y')}")
            include_title_slide = kwargs.get('include_title_slide', True)
            slides = kwargs.get('slides', [])

            slide_index = 0

            # Add title slide if requested
            if include_title_slide:
                title_payload = {
                    "index": slide_index,
                    "layoutInfo": {
                        "layoutType": "Title"
                    },
                    "title": title,
                    "subtitle": subtitle
                }

                # Create title slide
                requests.post(
                    f"https://graph.microsoft.com/v1.0/me/drive/items/{file_id}/workbook/presentation/slides",
                    headers=headers,
                    json=title_payload
                ).raise_for_status()

                slide_index += 1

            # Add agenda slide if requested
            if kwargs.get('include_agenda', False):
                agenda_title = kwargs.get('agenda_title', 'Agenda')

                # Generate agenda items
                agenda_items = kwargs.get('agenda_items', [])
                if kwargs.get('auto_generate_agenda', True):
                    agenda_items = [
                        slide.get('title', 'Untitled Slide') for slide in slides]

                # Build content as HTML
                agenda_html = "<ul>"
                for item in agenda_items:
                    agenda_html += f"<li>{item}</li>"
                agenda_html += "</ul>"

                agenda_payload = {
                    "index": slide_index,
                    "layoutInfo": {
                        "layoutType": "TitleAndContent"
                    },
                    "title": agenda_title,
                    "content": agenda_html
                }

                # Create agenda slide
                requests.post(
                    f"https://graph.microsoft.com/v1.0/me/drive/items/{file_id}/workbook/presentation/slides",
                    headers=headers,
                    json=agenda_payload
                ).raise_for_status()

                slide_index += 1

            # Add section slides at their specified positions (but adjusted for title/agenda slides)
            section_slides = kwargs.get('section_slides', [])
            section_positions = {}
            for section in section_slides:
                # Store position info for later use (need to adjust for any shifts)
                adj_position = section.get('position', 0) + slide_index
                section_positions[adj_position] = section

            # Add content slides
            for i, slide_data in enumerate(slides):
                slide_title = slide_data.get('title', '')
                content_items = slide_data.get('content', [])
                layout_name = slide_data.get(
                    'layout', 'title_and_content').lower()

                # Check if we need to insert a section slide first
                if slide_index in section_positions:
                    section_data = section_positions[slide_index]
                    section_payload = {
                        "index": slide_index,
                        "layoutInfo": {
                            "layoutType": "Section"
                        },
                        "title": section_data.get('title', 'New Section')
                    }

                    # Create section slide
                    requests.post(
                        f"https://graph.microsoft.com/v1.0/me/drive/items/{file_id}/workbook/presentation/slides",
                        headers=headers,
                        json=section_payload
                    ).raise_for_status()

                    slide_index += 1

                # Map layout name to Graph API layout type
                layout_type = "TitleAndContent"  # Default
                if layout_name == 'title':
                    layout_type = "Title"
                elif layout_name == 'section':
                    layout_type = "Section"
                elif layout_name == 'two_content':
                    layout_type = "TwoContent"
                elif layout_name == 'title_only':
                    layout_type = "TitleOnly"
                elif layout_name == 'blank':
                    layout_type = "Blank"
                elif layout_name == 'comparison' or layout_name == 'quote':
                    # These are custom layouts but we'll use the closest match
                    layout_type = "TwoContent" if layout_name == 'comparison' else "TitleAndContent"

                # Build content as HTML
                content_html = "<ul>"

                # Handle different layout types
                if layout_name in ['two_content', 'comparison']:
                    # Left content
                    left_content = slide_data.get('content_left', [])
                    content_html += "<div style='float: left; width: 45%;'><ul>"
                    for item in left_content:
                        is_bullet = item.strip().startswith('-') or item.strip().startswith('*')
                        is_numbered = bool(re.match(r'^\d+\.', item.strip()))

                        if is_bullet:
                            clean_item = item.strip()[1:].strip()
                            content_html += f"<li>{clean_item}</li>"
                        elif is_numbered:
                            clean_item = re.sub(
                                r'^\d+\.', '', item.strip()).strip()
                            content_html += f"<li>{clean_item}</li>"
                        else:
                            content_html += f"<p>{item}</p>"
                    content_html += "</ul></div>"

                    # Right content
                    right_content = slide_data.get('content_right', [])
                    content_html += "<div style='float: right; width: 45%;'><ul>"
                    for item in right_content:
                        is_bullet = item.strip().startswith('-') or item.strip().startswith('*')
                        is_numbered = bool(re.match(r'^\d+\.', item.strip()))

                        if is_bullet:
                            clean_item = item.strip()[1:].strip()
                            content_html += f"<li>{clean_item}</li>"
                        elif is_numbered:
                            clean_item = re.sub(
                                r'^\d+\.', '', item.strip()).strip()
                            content_html += f"<li>{clean_item}</li>"
                        else:
                            content_html += f"<p>{item}</p>"
                    content_html += "</ul></div>"

                elif layout_name == 'quote':
                    # Special formatting for quote layout
                    content_html = "<div style='font-style: italic; font-size: 24pt; text-align: center; margin: 50px;'>"
                    for item in content_items:
                        content_html += f"<p>\"{item}\"</p>"
                    content_html += "</div>"

                else:
                    # Standard content handling
                    for item in content_items:
                        is_bullet = item.strip().startswith('-') or item.strip().startswith('*')
                        is_numbered = bool(re.match(r'^\d+\.', item.strip()))

                        if is_bullet:
                            clean_item = item.strip()[1:].strip()
                            content_html += f"<li>{clean_item}</li>"
                        elif is_numbered:
                            clean_item = re.sub(
                                r'^\d+\.', '', item.strip()).strip()
                            content_html += f"<li>{clean_item}</li>"
                        else:
                            content_html += f"<p>{item}</p>"

                content_html += "</ul>"

                slide_payload = {
                    "index": slide_index,
                    "layoutInfo": {
                        "layoutType": layout_type
                    },
                    "title": slide_title,
                    "content": content_html
                }

                # Add subtitle if provided
                if slide_data.get('subtitle'):
                    slide_payload["subtitle"] = slide_data.get('subtitle')

                # Create content slide
                requests.post(
                    f"https://graph.microsoft.com/v1.0/me/drive/items/{file_id}/workbook/presentation/slides",
                    headers=headers,
                    json=slide_payload
                ).raise_for_status()

                slide_index += 1

            # Add thank you slide if requested
            if kwargs.get('thank_you_slide', False):
                thank_you_title = kwargs.get('thank_you_title', 'Thank You')
                thank_you_content = kwargs.get('thank_you_content', '')

                thank_you_payload = {
                    "index": slide_index,
                    "layoutInfo": {
                        "layoutType": "Title"
                    },
                    "title": thank_you_title
                }

                if thank_you_content:
                    thank_you_payload["subtitle"] = thank_you_content

                # Create thank you slide
                requests.post(
                    f"https://graph.microsoft.com/v1.0/me/drive/items/{file_id}/workbook/presentation/slides",
                    headers=headers,
                    json=thank_you_payload
                ).raise_for_status()

        except Exception as e:
            logging.error(f"Error adding slides with Graph API: {str(e)}")
            raise

    def _create_presentation_pptx(self, kwargs):
        """
        Creates a PowerPoint presentation using python-pptx (fallback method).

        Args:
            kwargs: All presentation parameters

        Returns:
            bytes: The binary presentation data
        """
        try:
            # Create new presentation
            prs = self.Presentation()

            # Extract parameters
            title = kwargs.get('title', 'Untitled Presentation')
            subtitle = kwargs.get(
                'subtitle', f"Created on {datetime.now().strftime('%B %d, %Y')}")
            author = kwargs.get('author', '')
            include_title_slide = kwargs.get('include_title_slide', True)
            slides_data = kwargs.get('slides', [])
            theme_colors = self._get_theme_colors(kwargs)
            font_name = kwargs.get('font', 'calibri')

            # Add title slide if requested
            if include_title_slide:
                title_slide_layout = prs.slide_layouts[0]  # Title slide layout
                slide = prs.slides.add_slide(title_slide_layout)

                # Set title
                if hasattr(slide.shapes, 'title') and slide.shapes.title:
                    title_shape = slide.shapes.title
                    title_shape.text = title
                    self._apply_text_style(
                        title_shape.text_frame, font_name, theme_colors, is_title=True)

                # Set subtitle
                subtitle_shape = None
                for shape in slide.placeholders:
                    if shape.placeholder_format.type == 2:  # Subtitle placeholder
                        subtitle_shape = shape
                        break

                if subtitle_shape:
                    subtitle_shape.text = subtitle
                    self._apply_text_style(
                        subtitle_shape.text_frame, font_name, theme_colors)

                # Add author if provided
                if author and subtitle_shape:
                    subtitle_shape.text += f"\n{author}"

                # Add image placeholder if specified
                if kwargs.get('title_slide_image'):
                    self._add_image_placeholder(slide, kwargs.get(
                        'title_slide_image'), theme_colors)

           # Add agenda slide if requested
            next_slide_idx = 1 if include_title_slide else 0
            if kwargs.get('include_agenda', False):
                agenda_title = kwargs.get('agenda_title', 'Agenda')

                # Create agenda slide
                # Title and content layout
                content_layout = prs.slide_layouts[1]
                agenda_slide = prs.slides.add_slide(content_layout)

                # Set title
                if hasattr(agenda_slide.shapes, 'title') and agenda_slide.shapes.title:
                    title_shape = agenda_slide.shapes.title
                    title_shape.text = agenda_title
                    self._apply_text_style(
                        title_shape.text_frame, font_name, theme_colors, is_title=True)

                # Find content placeholder
                content_placeholder = None
                for shape in agenda_slide.placeholders:
                    if shape.placeholder_format.type == 7:  # Content placeholder
                        content_placeholder = shape
                        break

                if content_placeholder:
                    tf = content_placeholder.text_frame
                    tf.clear()

                    # Generate agenda items
                    agenda_items = kwargs.get('agenda_items', [])
                    if kwargs.get('auto_generate_agenda', True):
                        agenda_items = [
                            slide.get('title', 'Untitled Slide') for slide in slides_data]

                    # Add items as bullet points
                    for i, item in enumerate(agenda_items):
                        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                        p.text = item
                        p.level = 0
                        p.font.size = Pt(18)
                        p.font.name = font_name
                        p.font.color.rgb = RGBColor(
                            *self._hex_to_rgb(theme_colors['text']))

                next_slide_idx += 1

            # Prepare section slides
            section_slides = kwargs.get('section_slides', [])
            section_positions = {}
            for section in section_slides:
                section_positions[section.get(
                    'position', 0) + next_slide_idx] = section

            # Track actual slide index for section insertion
            current_slide_idx = next_slide_idx

            # Add content slides
            for slide_data in slides_data:
                # Check if we need to insert a section slide first
                if current_slide_idx in section_positions:
                    section_data = section_positions[current_slide_idx]
                    section_slide = self._create_section_slide(
                        prs,
                        section_data.get('title', 'New Section'),
                        font_name,
                        theme_colors,
                        section_data.get('background_color')
                    )
                    current_slide_idx += 1

                # Create the content slide
                slide = self._create_content_slide(
                    prs, slide_data, font_name, theme_colors)
                current_slide_idx += 1

            # Add thank you slide if requested
            if kwargs.get('thank_you_slide', False):
                thank_you_title = kwargs.get('thank_you_title', 'Thank You')
                thank_you_content = kwargs.get('thank_you_content', '')

                # Create thank you slide using title slide layout
                title_slide_layout = prs.slide_layouts[0]
                slide = prs.slides.add_slide(title_slide_layout)

                # Set title
                if hasattr(slide.shapes, 'title') and slide.shapes.title:
                    title_shape = slide.shapes.title
                    title_shape.text = thank_you_title
                    self._apply_text_style(
                        title_shape.text_frame, font_name, theme_colors, is_title=True)

                # Set content as subtitle if provided
                if thank_you_content:
                    subtitle_shape = None
                    for shape in slide.placeholders:
                        if shape.placeholder_format.type == 2:  # Subtitle placeholder
                            subtitle_shape = shape
                            break

                    if subtitle_shape:
                        subtitle_shape.text = thank_you_content
                        self._apply_text_style(
                            subtitle_shape.text_frame, font_name, theme_colors)

            # Save to BytesIO and return binary data
            memory_file = BytesIO()
            prs.save(memory_file)
            memory_file.seek(0)
            return memory_file.getvalue()

        except Exception as e:
            logging.error(
                f"Error creating PowerPoint with python-pptx: {str(e)}")
            logging.error(traceback.format_exc())
            raise

        def _create_section_slide(self, prs, title, font_name, theme_colors, background_color=None):
            """
            Creates a section slide.

            Args:
                prs: Presentation object
                title: Section title
                font_name: Font to use
                theme_colors: Theme colors dictionary
                background_color: Optional background color override

            Returns:
                slide: The created slide
            """
            # Use section header layout or title only as fallback
            section_layout_idx = 2  # Typical index for section layout
            if section_layout_idx >= len(prs.slide_layouts):
                section_layout_idx = 5  # Title only layout as fallback

            slide_layout = prs.slide_layouts[section_layout_idx]
            slide = prs.slides.add_slide(slide_layout)

            # Set custom background if provided
            if background_color:
                background = slide.background
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(
                    *self._hex_to_rgb(background_color.lstrip('#')))
            else:
                # Use theme secondary color
                background = slide.background
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(
                    *self._hex_to_rgb(theme_colors['secondary']))

            # Set title
            if hasattr(slide.shapes, 'title') and slide.shapes.title:
                title_shape = slide.shapes.title
                title_shape.text = title

                # Special styling for section title
                tf = title_shape.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                p.text = title
                p.font.size = Pt(44)
                p.font.bold = True
                p.font.name = font_name
                # White text for contrast
                p.font.color.rgb = RGBColor(255, 255, 255)
                p.alignment = PP_ALIGN.CENTER

                # Center vertically
                title_shape.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

            return slide

    def _create_content_slide(self, prs, slide_data, font_name, theme_colors):
        """
        Creates a content slide based on slide data.

        Args:
            prs: Presentation object
            slide_data: Slide data dictionary
            font_name: Font to use
            theme_colors: Theme colors dictionary

        Returns:
            slide: The created slide
        """
        layout_name = slide_data.get('layout', 'title_and_content').lower()

        # Map layout name to index
        layout_idx = 1  # Default to title and content
        if layout_name == 'title':
            layout_idx = 0
        elif layout_name == 'section':
            layout_idx = 2
        elif layout_name == 'two_content':
            layout_idx = 3
        elif layout_name == 'title_only':
            layout_idx = 5
        elif layout_name == 'blank':
            layout_idx = 6
        elif layout_name == 'comparison':
            layout_idx = 3  # Use two content layout
        elif layout_name == 'quote':
            layout_idx = 5  # Use title only layout

        # Check if the layout index is valid
        if layout_idx >= len(prs.slide_layouts):
            layout_idx = 1  # Fallback to title and content

        # Create slide with selected layout
        slide_layout = prs.slide_layouts[layout_idx]
        slide = prs.slides.add_slide(slide_layout)

        # Set custom background if provided
        if slide_data.get('background_color'):
            background = slide.background
            fill = background.fill
            fill.solid()
            bg_color = slide_data.get('background_color').lstrip('#')
            fill.fore_color.rgb = RGBColor(*self._hex_to_rgb(bg_color))

        # Set slide title if available
        if hasattr(slide.shapes, 'title') and slide.shapes.title:
            title_shape = slide.shapes.title
            title_shape.text = slide_data.get('title', '')
            self._apply_text_style(
                title_shape.text_frame, font_name, theme_colors, is_title=True)

        # Add subtitle if provided and layout is title only or blank
        if slide_data.get('subtitle') and layout_name in ['title_only', 'blank']:
            # Add a text box for subtitle
            left = Inches(1)
            top = Inches(2)
            width = Inches(8)
            height = Inches(1)

            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.text = slide_data.get('subtitle')
            self._apply_text_style(tf, font_name, theme_colors)

        # Handle special layout: Quote
        if layout_name == 'quote':
            self._create_quote_slide(
                slide, slide_data, font_name, theme_colors)

        # Handle different content layouts
        elif layout_name in ['title_and_content', 'section']:
            self._add_standard_content(
                slide, slide_data, font_name, theme_colors)

        elif layout_name in ['two_content', 'comparison']:
            self._add_two_column_content(
                slide, slide_data, font_name, theme_colors)

        # Add image placeholder if specified
        if slide_data.get('image_placeholder'):
            self._add_image_placeholder(slide, slide_data.get(
                'image_placeholder'), theme_colors)

        # Add speaker notes if provided
        if slide_data.get('notes'):
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = slide_data.get('notes')

        return slide

    def _create_quote_slide(self, slide, slide_data, font_name, theme_colors):
        """
        Creates a quote slide with special formatting.

        Args:
            slide: Slide object
            slide_data: Slide data
            font_name: Font name
            theme_colors: Theme colors
        """
        content_items = slide_data.get('content', [])
        if not content_items:
            return

        # Create a text box for the quote
        left = Inches(1)
        top = Inches(2.5)
        width = Inches(8)
        height = Inches(3)

        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True

        # Add each quote line
        for i, quote in enumerate(content_items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f'"{quote}"'
            p.font.size = Pt(28)
            p.font.italic = True
            p.font.name = font_name
            p.font.color.rgb = RGBColor(
                *self._hex_to_rgb(theme_colors['primary']))
            p.alignment = PP_ALIGN.CENTER

    def _add_standard_content(self, slide, slide_data, font_name, theme_colors):
        """
        Adds content to a standard slide.

        Args:
            slide: Slide object
            slide_data: Slide data dictionary
            font_name: Font name
            theme_colors: Theme colors dictionary
        """
        content_items = slide_data.get('content', [])
        if not content_items:
            return

        # Find content placeholder
        content_placeholder = None
        for shape in slide.placeholders:
            if shape.placeholder_format.type == 7:  # Content placeholder
                content_placeholder = shape
                break

        if content_placeholder:
            tf = content_placeholder.text_frame
            tf.clear()  # Clear any default text

            # Process content items
            for i, text in enumerate(content_items):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()

                # Handle bullet points
                is_bullet = text.strip().startswith('-') or text.strip().startswith('*')
                is_numbered = bool(re.match(r'^\d+\.', text.strip()))

                if is_bullet:
                    # Remove bullet character and set as bullet point
                    p.text = text.strip()[1:].strip()
                    p.level = 0  # First level bullet
                elif is_numbered:
                    # Remove number prefix
                    p.text = re.sub(r'^\d+\.', '', text.strip()).strip()
                    p.level = 0  # First level bullet
                else:
                    p.text = text
                    p.level = 0  # Regular paragraph

                # Apply style
                p.font.size = Pt(18)
                p.font.name = font_name
                p.font.color.rgb = RGBColor(
                    *self._hex_to_rgb(theme_colors['text']))

    def _add_two_column_content(self, slide, slide_data, font_name, theme_colors):
        """
        Adds content in two columns.

        Args:
            slide: Slide object
            slide_data: Slide data dictionary
            font_name: Font name
            theme_colors: Theme colors dictionary
        """
        # Get content
        left_content = slide_data.get('content_left', [])
        right_content = slide_data.get('content_right', [])

        # If standard content is provided but not left/right, split it
        if not left_content and not right_content:
            content = slide_data.get('content', [])
            mid_point = len(content) // 2
            left_content = content[:mid_point]
            right_content = content[mid_point:]

        # Find placeholders
        left_placeholder = right_placeholder = None

        for shape in slide.placeholders:
            # Left content placeholder (usually index 1)
            if shape.placeholder_format.idx == 1:
                left_placeholder = shape
            # Right content placeholder (usually index 2)
            elif shape.placeholder_format.idx == 2:
                right_placeholder = shape

        # Add left column content
        if left_placeholder and left_content:
            tf = left_placeholder.text_frame
            tf.clear()

            for i, text in enumerate(left_content):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()

                is_bullet = text.strip().startswith('-') or text.strip().startswith('*')
                is_numbered = bool(re.match(r'^\d+\.', text.strip()))

                if is_bullet:
                    p.text = text.strip()[1:].strip()
                    p.level = 0
                elif is_numbered:
                    p.text = re.sub(r'^\d+\.', '', text.strip()).strip()
                    p.level = 0
                else:
                    p.text = text
                    p.level = 0

                # Apply style
                p.font.size = Pt(16)
                p.font.name = font_name
                p.font.color.rgb = RGBColor(
                    *self._hex_to_rgb(theme_colors['text']))

        # Add right column content
        if right_placeholder and right_content:
            tf = right_placeholder.text_frame
            tf.clear()

            for i, text in enumerate(right_content):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()

                is_bullet = text.strip().startswith('-') or text.strip().startswith('*')
                is_numbered = bool(re.match(r'^\d+\.', text.strip()))

                if is_bullet:
                    p.text = text.strip()[1:].strip()
                    p.level = 0
                elif is_numbered:
                    p.text = re.sub(r'^\d+\.', '', text.strip()).strip()
                    p.level = 0
                else:
                    p.text = text
                    p.level = 0

                # Apply style
                p.font.size = Pt(16)
                p.font.name = font_name
                p.font.color.rgb = RGBColor(
                    *self._hex_to_rgb(theme_colors['text']))

    def _add_image_placeholder(self, slide, image_spec, theme_colors):
        """
        Adds an image placeholder with a colored shape.

        Args:
            slide: Slide object
            image_spec: Image specification string (format: image:width:height:position)
            theme_colors: Theme colors dictionary
        """
        try:
            # Parse image spec
            parts = image_spec.split(':')

            # Default values
            width = Inches(4)
            height = Inches(3)
            position = 'center'

            # Override with provided values
            if len(parts) >= 3:
                try:
                    width = Inches(float(parts[1]))
                    height = Inches(float(parts[2]))
                except ValueError:
                    pass

            if len(parts) >= 4:
                position = parts[3].lower()

            # Calculate position
            slide_width = Inches(10)  # Standard slide width
            slide_height = Inches(7.5)  # Standard slide height

            # Default position (center)
            left = (slide_width - width) / 2
            top = (slide_height - height) / 2

            # Adjust based on position keyword
            if position == 'left':
                left = Inches(0.5)
            elif position == 'right':
                left = slide_width - width - Inches(0.5)
            elif position == 'top':
                top = Inches(1.5)
            elif position == 'bottom':
                top = slide_height - height - Inches(0.5)

            # Create a colored rectangle as placeholder
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                left, top, width, height
            )

            # Style the shape
            fill = shape.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(
                *self._hex_to_rgb(theme_colors['accent1']))

            # Add text label
            tf = shape.text_frame
            tf.text = f"Placeholder for {parts[0]}"
            tf.word_wrap = True

            # Center text
            for paragraph in tf.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                paragraph.font.color.rgb = RGBColor(
                    255, 255, 255)  # White text

            # Center vertically
            tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

        except Exception as e:
            logging.warning(f"Error adding image placeholder: {str(e)}")

    def _apply_text_style(self, text_frame, font_name, theme_colors, is_title=False):
        """
        Applies consistent text styling.

        Args:
            text_frame: Text frame to style
            font_name: Font name to use
            theme_colors: Theme colors
            is_title: Whether this is a title text frame
        """
        for paragraph in text_frame.paragraphs:
            if is_title:
                paragraph.font.size = Pt(32)
                paragraph.font.bold = True
                paragraph.font.color.rgb = RGBColor(
                    *self._hex_to_rgb(theme_colors['primary']))
            else:
                paragraph.font.size = Pt(18)
                paragraph.font.color.rgb = RGBColor(
                    *self._hex_to_rgb(theme_colors['text']))

            paragraph.font.name = font_name

    def _get_access_token(self):
        """
        Gets an access token for Microsoft Graph API.

        Returns:
            str: Access token or None if failed
        """
        try:
            token_url = f"https://login.microsoftonline.com/{self.tenant_id}/oauth2/v2.0/token"

            token_data = {
                'grant_type': 'client_credentials',
                'client_id': self.client_id,
                'client_secret': self.client_secret,
                'scope': 'https://graph.microsoft.com/.default'
            }

            response = requests.post(token_url, data=token_data)
            response.raise_for_status()

            token_info = response.json()
            return token_info.get('access_token')

        except Exception as e:
            logging.error(f"Error getting Graph API access token: {str(e)}")
            return None

    def _get_as_base64(self, presentation_data, title):
        """
        Converts the presentation to base64 format.

        Args:
            presentation_data (bytes): Binary presentation data
            title (str): The title of the presentation

        Returns:
            dict: Base64-encoded presentation data
        """
        try:
            # Encode binary data to base64
            base64_content = base64.b64encode(
                presentation_data).decode('utf-8')

            # Create sanitized filename
            sanitized_title = self._sanitize_filename(title)
            filename = f"{sanitized_title}.pptx"

            return {
                "status": "success",
                "message": "PowerPoint presentation created successfully",
                "filename": filename,
                "content_type": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                "base64_data": base64_content
            }

        except Exception as e:
            logging.error(f"Error converting to base64: {str(e)}")
            logging.error(traceback.format_exc())
            return {"status": "error", "message": f"Failed to convert presentation to base64: {str(e)}"}

    def _save_to_local_file(self, presentation_data, title, file_path, filename_prefix=None):
        """
        Saves the presentation to a local file.

        Args:
            presentation_data (bytes): Binary presentation data
            title (str): The title of the presentation
            file_path (str): Path to save the file
            filename_prefix (str): Optional prefix for filename

        Returns:
            dict: Result of the save operation
        """
        try:
            # Generate filename if not specified
            if not file_path:
                sanitized_title = self._sanitize_filename(title)
                timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')

                if filename_prefix:
                    sanitized_title = f"{filename_prefix}_{sanitized_title}"

                file_path = os.path.join(
                    tempfile.gettempdir(),
                    f"{sanitized_title}_{timestamp}.pptx"
                )

            # Ensure directory exists
            os.makedirs(os.path.dirname(
                os.path.abspath(file_path)), exist_ok=True)

            # Write binary data to file
            with open(file_path, 'wb') as f:
                f.write(presentation_data)

            return {
                "status": "success",
                "message": f"PowerPoint presentation saved to {file_path}",
                "file_path": file_path
            }

        except Exception as e:
            logging.error(f"Error saving to file: {str(e)}")
            logging.error(traceback.format_exc())
            return {"status": "error", "message": f"Failed to save presentation to file: {str(e)}"}

    def _save_to_azure(self, presentation_data, title, azure_directory=None, user_guid=None, filename_prefix=None):
        """
        Saves the presentation to Azure File Storage using a standardized location.

        Args:
            presentation_data (bytes): Binary presentation data
            title (str): The title of the presentation
            azure_directory (str): Optional subdirectory (will be created under standard path)
            user_guid (str): User GUID for user-specific storage
            filename_prefix (str): Optional prefix for filename

        Returns:
            dict: Result of the save operation with detailed path information
        """
        try:
            # Generate filename with timestamp
            sanitized_title = self._sanitize_filename(title)
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')

            if filename_prefix:
                sanitized_title = f"{filename_prefix}_{sanitized_title}"

            filename = f"{sanitized_title}_{timestamp}.pptx"

            # Use standardized storage paths
            # Root directory is always "powerpoint_presentations"
            storage_root = "powerpoint_presentations"

            if user_guid:
                # User-specific storage
                # Format: powerpoint_presentations/users/{user_guid}/{optional_subdir}/
                self.storage_manager.set_memory_context(user_guid)
                storage_dir = f"{storage_root}/users/{user_guid}"
            else:
                # Shared storage
                # Format: powerpoint_presentations/shared/{optional_subdir}/
                storage_dir = f"{storage_root}/shared"

            # Add subdirectory if specified
            if azure_directory and azure_directory.strip():
                # Remove any leading/trailing slashes for consistency
                clean_subdir = azure_directory.strip().strip('/')
                storage_dir = f"{storage_dir}/{clean_subdir}"

            # Create date-based subdirectory for better organization
            date_subdir = datetime.now().strftime('%Y-%m')
            storage_dir = f"{storage_dir}/{date_subdir}"

            # Ensure directory exists
            self.storage_manager.ensure_directory_exists(storage_dir)

            # Upload to Azure File Storage
            success = self.storage_manager.write_file(
                storage_dir,
                filename,
                presentation_data
            )

            # Construct a user-friendly path for display
            if user_guid:
                display_path = f"Your PowerPoint Presentations > {date_subdir}"
                if azure_directory:
                    display_path = f"Your PowerPoint Presentations > {azure_directory} > {date_subdir}"
            else:
                display_path = f"Shared PowerPoint Presentations > {date_subdir}"
                if azure_directory:
                    display_path = f"Shared PowerPoint Presentations > {azure_directory} > {date_subdir}"

            if success:
                return {
                    "status": "success",
                    "message": f"PowerPoint presentation saved successfully to: {display_path}",
                    "storage_path": f"{storage_dir}/{filename}",
                    "display_path": display_path,
                    "filename": filename,
                    # For system use
                    "technical_path": f"{storage_dir}/{filename}"
                }
            else:
                return {
                    "status": "error",
                    "message": "Failed to save PowerPoint presentation to Azure File Storage"
                }

        except Exception as e:
            logging.error(f"Error saving to Azure: {str(e)}")
            logging.error(traceback.format_exc())
            return {
                "status": "error",
                "message": f"Error saving presentation to Azure: {str(e)}"
            }

    def _sanitize_filename(self, filename):
        """
        Sanitizes the filename to be valid.

        Args:
            filename (str): Original filename

        Returns:
            str: Sanitized filename
        """
        # Replace invalid characters with underscores
        invalid_chars = '<>:"/\\|?*'
        for char in invalid_chars:
            filename = filename.replace(char, '_')

        # Trim to reasonable length
        if len(filename) > 100:
            filename = filename[:97] + '...'

        return filename
